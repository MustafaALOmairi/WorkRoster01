"""Build WorkRoster_Presentation.pptx (English, formal, simple)."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from PIL import Image
import os

# Brand
ACCENT = RGBColor(0xF5, 0xA6, 0x23)      # warm orange (app accent)
DARK = RGBColor(0x0D, 0x11, 0x17)        # near-black
INK = RGBColor(0x1F, 0x29, 0x37)         # body text
MUTED = RGBColor(0x6B, 0x72, 0x80)       # caption
LIGHT = RGBColor(0xF7, 0xF8, 0xFA)       # card bg
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
RULE = RGBColor(0xE5, 0xE7, 0xEB)

W, H = Inches(13.333), Inches(7.5)

prs = Presentation()
prs.slide_width = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]


def add_rect(slide, x, y, w, h, fill, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
    shp.shadow.inherit = False
    return shp


def add_text(slide, x, y, w, h, text, *, size=18, bold=False, color=INK,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Calibri"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    lines = text.split("\n")
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = ln
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
    return tb


def header(slide, title, kicker=None):
    # Top accent bar
    add_rect(slide, 0, 0, W, Inches(0.18), ACCENT)
    # Brand mark (top-left)
    add_text(slide, Inches(0.6), Inches(0.32), Inches(4), Inches(0.4),
             "WorkRoster", size=14, bold=True, color=ACCENT)
    if kicker:
        add_text(slide, Inches(0.6), Inches(0.85), Inches(8), Inches(0.4),
                 kicker.upper(), size=11, bold=True, color=MUTED)
    add_text(slide, Inches(0.6), Inches(1.2), Inches(12), Inches(0.85),
             title, size=34, bold=True, color=DARK)
    # underline
    add_rect(slide, Inches(0.6), Inches(2.05), Inches(0.7), Emu(38100), ACCENT)
    # page footer
    add_text(slide, Inches(0.6), Inches(7.05), Inches(8), Inches(0.3),
             "WorkRoster  ·  Smart Scheduling for Shift Workers",
             size=9, color=MUTED)


def add_image_card(slide, img_path, x, y, w, h):
    """Place a shadowed white card with a fitted image inside."""
    # card
    card = add_rect(slide, x, y, w, h, WHITE, line=RULE)
    # fit image
    if not os.path.exists(img_path):
        add_text(slide, x, y, w, h, "[image missing]", size=12,
                 color=MUTED, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        return
    pad = Inches(0.15)
    iw, ih = w - 2 * pad, h - 2 * pad
    with Image.open(img_path) as im:
        ar = im.width / im.height
    box_ar = iw / ih
    if ar > box_ar:
        pic_w = iw; pic_h = int(iw / ar)
    else:
        pic_h = ih; pic_w = int(ih * ar)
    px = x + (w - pic_w) // 2
    py = y + (h - pic_h) // 2
    slide.shapes.add_picture(img_path, px, py, pic_w, pic_h)


def bullet_list(slide, x, y, w, h, items, size=16, gap=0.35):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap * 14)
        r = p.add_run()
        r.text = "•  " + it
        r.font.name = "Calibri"
        r.font.size = Pt(size)
        r.font.color.rgb = INK


# ---------------- SLIDE 1 — TITLE ----------------
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, W, H, DARK)
add_rect(s, 0, Inches(3.2), W, Inches(0.05), ACCENT)
add_text(s, Inches(0.8), Inches(2.0), Inches(11.7), Inches(0.6),
         "WORKROSTER", size=14, bold=True, color=ACCENT)
add_text(s, Inches(0.8), Inches(2.4), Inches(11.7), Inches(1.2),
         "A Smart Calendar for Shift Workers", size=44, bold=True, color=WHITE)
add_text(s, Inches(0.8), Inches(3.5), Inches(11.7), Inches(0.6),
         "Project Report  ·  Version 1.0.1  ·  2026", size=16, color=RGBColor(0xC8, 0xCD, 0xD3))
add_text(s, Inches(0.8), Inches(6.7), Inches(11.7), Inches(0.4),
         "Mobile  ·  iOS  ·  Android  ·  Web", size=12, color=MUTED)


# ---------------- SLIDE 2 — THE PROBLEM ----------------
s = prs.slides.add_slide(BLANK)
header(s, "The Problem Before WorkRoster", kicker="Section 01")
add_text(s, Inches(0.6), Inches(2.4), Inches(12), Inches(0.6),
         "Shift workers had no proper tool to organize their schedules.",
         size=20, bold=True, color=DARK)
bullet_list(s, Inches(0.7), Inches(3.2), Inches(12), Inches(3.5), [
    "Employees struggled to organize rotating shifts and track holidays.",
    "Schedules were managed manually using paper rosters or spreadsheets.",
    "Information was unclear, scattered, and easily out of date.",
    "Confusion in shift distribution led to missed days and overlap.",
    "No simple way to share holiday calendars between team members.",
], size=17)


# ---------------- SLIDE 3 — THE SOLUTION ----------------
s = prs.slides.add_slide(BLANK)
header(s, "Our Solution", kicker="Section 02")
add_text(s, Inches(0.6), Inches(2.4), Inches(12), Inches(0.6),
         "WorkRoster turns shift management into a smart, automated system.",
         size=20, bold=True, color=DARK)
# Two-column highlight cards
def hcard(x, title, body):
    add_rect(s, x, Inches(3.2), Inches(5.9), Inches(2.9), LIGHT)
    add_rect(s, x, Inches(3.2), Inches(0.12), Inches(2.9), ACCENT)
    add_text(s, x + Inches(0.35), Inches(3.4), Inches(5.4), Inches(0.5),
             title, size=18, bold=True, color=DARK)
    add_text(s, x + Inches(0.35), Inches(3.95), Inches(5.4), Inches(2.0),
             body, size=14, color=INK)

hcard(Inches(0.6),
      "Set up once, run forever",
      "Configure your shift pattern and start date a single time. "
      "WorkRoster keeps the calendar up to date automatically — "
      "every day, every month, every year.")
hcard(Inches(6.85),
      "Everything in one place",
      "Shifts, holidays, notes, reminders, and team holiday sharing — "
      "all unified inside a clean, bilingual mobile app that works "
      "online or fully offline.")


# ---------------- SLIDE 4 — HOW IT WORKS ----------------
s = prs.slides.add_slide(BLANK)
header(s, "How the App Works", kicker="Section 03")
# Three step cards
def step(x, num, title, body):
    add_rect(s, x, Inches(3.0), Inches(3.95), Inches(3.5), WHITE, line=RULE)
    # circle
    c = s.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.3), Inches(3.2),
                           Inches(0.7), Inches(0.7))
    c.fill.solid(); c.fill.fore_color.rgb = ACCENT
    c.line.fill.background()
    tf = c.text_frame; tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = num
    r.font.bold = True; r.font.size = Pt(20); r.font.color.rgb = WHITE; r.font.name = "Calibri"
    add_text(s, x + Inches(0.3), Inches(4.05), Inches(3.45), Inches(0.5),
             title, size=16, bold=True, color=DARK)
    add_text(s, x + Inches(0.3), Inches(4.55), Inches(3.45), Inches(1.8),
             body, size=13, color=INK)

step(Inches(0.6), "1", "Enter your shift schedule",
     "Add your rotation pattern (morning / evening / night / rest) "
     "and choose your start date.")
step(Inches(4.7), "2", "The app fills the calendar",
     "Daily and monthly views are generated automatically and updated "
     "as time passes.")
step(Inches(8.8), "3", "Receive smart reminders",
     "Get alerts for upcoming holidays and shift changes — never miss "
     "an important day.")


# ---------------- SLIDE 5 — FEATURES OVERVIEW ----------------
s = prs.slides.add_slide(BLANK)
header(s, "Key Features", kicker="Section 04")
features = [
    ("Shift Organization", "Morning, evening, night and rest shifts — color coded."),
    ("Daily & Monthly View", "See your schedule at a glance, any month."),
    ("Official Holidays", "Add and manage public holidays."),
    ("Custom Annual Leave", "Personalize your own vacation days."),
    ("Daily Notes", "Attach notes to any day."),
    ("Smart Reminders", "Get notified before key dates."),
    ("Simple Interface", "Clean design, fast to learn."),
    ("Custom Themes", "Light, dark, and AI-generated themes."),
    ("Powerful Search", "Find any day, any range, any holiday."),
]
# 3x3 grid
gx0, gy0 = Inches(0.6), Inches(2.5)
cw, ch = Inches(4.05), Inches(1.45)
gap_x, gap_y = Inches(0.1), Inches(0.15)
for i, (t, d) in enumerate(features):
    row, col = i // 3, i % 3
    x = gx0 + col * (cw + gap_x)
    y = gy0 + row * (ch + gap_y)
    add_rect(s, x, y, cw, ch, LIGHT)
    add_rect(s, x, y, Inches(0.08), ch, ACCENT)
    add_text(s, x + Inches(0.25), y + Inches(0.15), cw - Inches(0.4), Inches(0.4),
             t, size=14, bold=True, color=DARK)
    add_text(s, x + Inches(0.25), y + Inches(0.55), cw - Inches(0.4), Inches(0.85),
             d, size=11, color=INK)


# ---------------- SLIDE 6 — CALENDAR (image) ----------------
s = prs.slides.add_slide(BLANK)
header(s, "Feature in Focus — The Calendar", kicker="Feature 01")
add_image_card(s, "attached_assets/deck/01_calendar.jpg",
               Inches(0.6), Inches(2.4), Inches(7.4), Inches(4.4))
add_text(s, Inches(8.3), Inches(2.5), Inches(4.6), Inches(0.5),
         "Color-coded daily view", size=16, bold=True, color=DARK)
bullet_list(s, Inches(8.3), Inches(3.0), Inches(4.6), Inches(3.5), [
    "Each day is auto-filled with the correct shift.",
    "Tap any day to see details, add notes or set a reminder.",
    "Indicators show where notes or holidays exist.",
    "Adjustable day-number font size for accessibility.",
], size=12, gap=0.4)


# ---------------- SLIDE 7 — CUSTOMIZE (image) ----------------
s = prs.slides.add_slide(BLANK)
header(s, "Feature in Focus — Customize", kicker="Feature 02")
add_image_card(s, "attached_assets/deck/02_customize.jpg",
               Inches(0.6), Inches(2.4), Inches(7.4), Inches(4.4))
add_text(s, Inches(8.3), Inches(2.5), Inches(4.6), Inches(0.5),
         "Make the schedule yours", size=16, bold=True, color=DARK)
bullet_list(s, Inches(8.3), Inches(3.0), Inches(4.6), Inches(3.5), [
    "Set your shift start date.",
    "Choose a Quick Setup preset or build a custom rotation.",
    "Adjust working hours per shift type.",
    "Personalize the color of every shift category.",
], size=12, gap=0.4)


# ---------------- SLIDE 8 — SEARCH (image) ----------------
s = prs.slides.add_slide(BLANK)
header(s, "Feature in Focus — Search", kicker="Feature 03")
add_image_card(s, "attached_assets/deck/03_search.jpg",
               Inches(0.6), Inches(2.4), Inches(7.4), Inches(4.4))
add_text(s, Inches(8.3), Inches(2.5), Inches(4.6), Inches(0.5),
         "Find any day instantly", size=16, bold=True, color=DARK)
bullet_list(s, Inches(8.3), Inches(3.0), Inches(4.6), Inches(3.5), [
    "Search by a single date — see if it is work or off.",
    "Search a date range and view every shift inside it.",
    "Browse all upcoming holidays in one list.",
    "Tap any result to jump straight to that day.",
], size=12, gap=0.4)


# ---------------- SLIDE 9 — THEMES (image) ----------------
s = prs.slides.add_slide(BLANK)
header(s, "Feature in Focus — Themes", kicker="Feature 04")
add_image_card(s, "attached_assets/deck/04_themes.jpg",
               Inches(0.6), Inches(2.4), Inches(7.4), Inches(4.4))
add_text(s, Inches(8.3), Inches(2.5), Inches(4.6), Inches(0.5),
         "Make it look the way you like", size=16, bold=True, color=DARK)
bullet_list(s, Inches(8.3), Inches(3.0), Inches(4.6), Inches(3.5), [
    "Three premium built-in themes.",
    "Generate unlimited AI themes — describe a vibe, get a palette.",
    "Light and dark mode supported throughout.",
    "Themes apply instantly to the whole calendar.",
], size=12, gap=0.4)


# ---------------- SLIDE 10 — SETTINGS (image) ----------------
s = prs.slides.add_slide(BLANK)
header(s, "Feature in Focus — Settings", kicker="Feature 05")
add_image_card(s, "attached_assets/deck/05_settings.jpg",
               Inches(0.6), Inches(2.4), Inches(7.4), Inches(4.4))
add_text(s, Inches(8.3), Inches(2.5), Inches(4.6), Inches(0.5),
         "Personal & accessibility controls", size=16, bold=True, color=DARK)
bullet_list(s, Inches(8.3), Inches(3.0), Inches(4.6), Inches(3.5), [
    "Switch between Arabic and English.",
    "Light or dark appearance.",
    "Calendar number font size — Small to X-Large.",
    "Enable or mute interactive sounds.",
], size=12, gap=0.4)


# ---------------- SLIDE 11 — HOW TO USE ----------------
s = prs.slides.add_slide(BLANK)
header(s, "How to Use the App", kicker="Section 05")
steps = [
    ("01", "Open the app",
     "No signup required. The app is fully usable offline."),
    ("02", "Configure your shift",
     "Open Customize, pick a Quick Setup preset or build your own pattern, "
     "and set your start date."),
    ("03", "View your calendar",
     "Return to the home screen — your shifts are filled in for every month, "
     "past and future."),
    ("04", "Add notes & reminders",
     "Tap any day to view its shift, add personal notes, or set reminders."),
    ("05", "Sync across devices  (optional)",
     "Create a free account in Account to back up data and sync between phones."),
]
y = Inches(2.45)
for num, title, body in steps:
    # number circle
    c = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.7), y,
                           Inches(0.55), Inches(0.55))
    c.fill.solid(); c.fill.fore_color.rgb = ACCENT
    c.line.fill.background()
    tf = c.text_frame; tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = num
    r.font.bold = True; r.font.size = Pt(11); r.font.color.rgb = WHITE; r.font.name = "Calibri"
    add_text(s, Inches(1.5), y, Inches(11), Inches(0.4),
             title, size=15, bold=True, color=DARK)
    add_text(s, Inches(1.5), y + Inches(0.35), Inches(11), Inches(0.5),
             body, size=12, color=INK)
    y += Inches(0.85)


# ---------------- SLIDE 12 — ROADMAP ----------------
s = prs.slides.add_slide(BLANK)
header(s, "What's Next — Future Roadmap", kicker="Section 06")
add_text(s, Inches(0.6), Inches(2.4), Inches(12), Inches(0.6),
         "Coming soon: Team Lead module.", size=20, bold=True, color=DARK)
add_text(s, Inches(0.6), Inches(3.0), Inches(12), Inches(0.5),
         "Built for managers and team leaders, with the same simple experience.",
         size=13, color=MUTED)
roadmap = [
    ("View team schedules", "See the full roster of every team member at a glance."),
    ("Plan in advance", "Adjust shifts and balance coverage weeks ahead."),
    ("Coverage alerts", "Receive automatic warnings when a shift is understaffed."),
    ("Monthly reports", "Generate clear monthly attendance and shift reports."),
]
y = Inches(3.7)
for title, body in roadmap:
    add_rect(s, Inches(0.6), y, Inches(12.1), Inches(0.7), LIGHT)
    add_rect(s, Inches(0.6), y, Inches(0.1), Inches(0.7), ACCENT)
    add_text(s, Inches(0.85), y + Inches(0.1), Inches(4), Inches(0.5),
             title, size=14, bold=True, color=DARK)
    add_text(s, Inches(5.1), y + Inches(0.12), Inches(7.4), Inches(0.5),
             body, size=12, color=INK)
    y += Inches(0.78)


# ---------------- SLIDE 13 — SUMMARY ----------------
s = prs.slides.add_slide(BLANK)
header(s, "Summary", kicker="Section 07")
add_text(s, Inches(0.6), Inches(2.6), Inches(12), Inches(1.0),
         "WorkRoster turns shift management from a complicated task\ninto a smart, automated system.",
         size=24, bold=True, color=DARK)
# three pillars
def pillar(x, title, body):
    add_rect(s, x, Inches(4.5), Inches(3.95), Inches(2.0), LIGHT)
    add_text(s, x + Inches(0.25), Inches(4.7), Inches(3.7), Inches(0.5),
             title, size=15, bold=True, color=ACCENT)
    add_text(s, x + Inches(0.25), Inches(5.2), Inches(3.7), Inches(1.4),
             body, size=12, color=INK)
pillar(Inches(0.6), "Boosts efficiency",
       "Removes manual scheduling and reduces day-to-day errors.")
pillar(Inches(4.7), "Supports planning",
       "Lets workers and teams plan ahead with full visibility.")
pillar(Inches(8.8), "Stays simple",
       "A clean, bilingual interface anyone can use in minutes.")


# ---------------- SLIDE 14 — AVAILABILITY ----------------
s = prs.slides.add_slide(BLANK)
header(s, "Availability & Reach", kicker="Section 08")
add_text(s, Inches(0.6), Inches(2.5), Inches(12), Inches(0.6),
         "Ready to download on every major platform.",
         size=20, bold=True, color=DARK)
def store(x, name, sub):
    add_rect(s, x, Inches(3.5), Inches(3.95), Inches(2.4), WHITE, line=RULE)
    add_rect(s, x, Inches(3.5), Inches(3.95), Inches(0.5), DARK)
    add_text(s, x, Inches(3.6), Inches(3.95), Inches(0.4),
             name, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.3), Inches(4.2), Inches(3.4), Inches(1.6),
             sub, size=12, color=INK)
store(Inches(0.6), "Apple App Store",
      "iPhone & iPad. Submission in progress for v1.0.1.")
store(Inches(4.7), "Google Play Store",
      "Android phones and tablets.")
store(Inches(8.8), "Web App",
      "Run instantly in any modern browser — no install needed.")
add_text(s, Inches(0.6), Inches(6.2), Inches(12), Inches(0.6),
         "Already field-tested across multiple shift systems and companies — proving its flexibility and reliability.",
         size=12, color=MUTED, align=PP_ALIGN.CENTER)


# ---------------- SLIDE 15 — THANK YOU ----------------
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, W, H, DARK)
add_rect(s, 0, Inches(4.0), W, Inches(0.05), ACCENT)
add_text(s, Inches(0.8), Inches(2.8), Inches(11.7), Inches(1.2),
         "Thank You", size=60, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.8), Inches(4.2), Inches(11.7), Inches(0.6),
         "WorkRoster — Your shifts, simplified.",
         size=20, color=ACCENT, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.8), Inches(5.0), Inches(11.7), Inches(0.5),
         "Questions? We're happy to walk you through any part of the app.",
         size=14, color=RGBColor(0xC8, 0xCD, 0xD3), align=PP_ALIGN.CENTER)


out = "WorkRoster_Presentation.pptx"
prs.save(out)
print("OK", out, os.path.getsize(out))
